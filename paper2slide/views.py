from django.http import HttpResponse
import collections
import collections.abc
import json
from pptx import Presentation

# TODO: Jongwon


def index(request):
    return HttpResponse("P2S index.")

# TODO: Heejae


def pdf_to_text(pdf_file):
    """
    Divide given paper (pdf file) into sections,
    (e.g. Intorduction, Background, Evaluation ...)
    and make it to text format.

    :param pdf_file: input paper pdf file
    :return: proper format for describing papaer
    (The following format is just an example. You can define better format.)
    (e.g. [
        {
            'title': 'Introduction',
            'body': 'Operating system kernel schedulers are ...'
        },
        {
            'title': 'Overview of CFS and ULE',
            'body': [
                {
                    'title': 'Linux CFS',
                    'body': 'Per-core scheduling: Linuxs CFS implements ...'
                },
                {
                    'title': 'FreeBSD ULE',
                    'body': 'Per-core scheduling: ULE uses two runqueues ...'
                }
            ]
        },
        ...
    ])
    """
    pass

# TODO: Inseo


slide_layout = {"Title": 0,
                "Title and Content": 1,
                "Section Header": 2,
                "Two Content": 3,
                "Comparison": 4,
                "Title Only": 5,
                "Blank": 6,
                "Contnet with Caption": 7,
                "Picture with Caption": 8}


def generate_slide(paper_summary):
    paper_summary = json.loads(paper_summary)
    prs = Presentation()  # 새 슬라이드 생성

    name = paper_summary["UserInfo"]["UserID"] + \
        "_" + paper_summary["UserInfo"]["RequestID"]

    for slide_info in paper_summary["slideInfos"]:
        layout = prs.slide_layouts[slide_layout[slide_info["layout"]]]
        slide = prs.slides.add_slide(layout)

        for placeholder_info in slide_info["placeholder"]:

            # print("\n placeholder_info: ", placeholder_info)
            # print("placeholder_info type: ", type(placeholder_info), "\n")

            placeholder_type = placeholder_info["type"]

            if (placeholder_type == "title"):
                slide.placeholders[placeholder_info["number"]
                                   ].text = placeholder_info["content"][0]["item"]

            elif (placeholder_type == "subtitle"):
                slide.placeholders[placeholder_info["number"]
                                   ].text = placeholder_info["content"][0]["item"]

            elif (placeholder_type == 'contents placeholder'):
                body = slide.shapes.placeholders[placeholder_info["number"]]
                tf = body.text_frame
                for content in placeholder_info["content"]:
                    # print("\n content: ", content, "\n")
                    if (content["type"] == "text"):
                        p = tf.add_paragraph()
                        p.text = content["item"]
                        p.level = content["level"]

            elif (placeholder_type == "text placeholder"):
                slide.placeholders[placeholder_info["number"]
                                   ].text = placeholder_info["content"][0]["item"]

            elif (placeholder_type == "picture placeholder"):
                # print("picture name: ", placeholder_info["content"][0]["item"])
                slide.placeholders[placeholder_info["number"]].insert_picture(
                    placeholder_info["content"][0]["item"])

    prs.save(name+'.pptx')
